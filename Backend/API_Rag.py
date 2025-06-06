import os
import chromadb
import fastapi
from fastapi import FastAPI, UploadFile, File, Response, Form, HTTPException, Depends, status, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse, StreamingResponse
from pydantic import BaseModel
from sentence_transformers import SentenceTransformer
import google.generativeai as genai
import base64
from PIL import Image
import io
import uuid
from datetime import datetime
import time
import pandas as pd
from openpyxl import Workbook
from openpyxl.drawing.image import Image as XLImage
from openpyxl.utils import get_column_letter
import shutil
import json
import openpyxl.styles
from fastapi.staticfiles import StaticFiles
from chromadb import PersistentClient
from chromadb.config import Settings

# 🔐 Đặt API key của Gemini từ biến môi trường thay vì hardcode (bảo mật hơn)
GENMINI_API_KEY = "AIzaSyAqX5bkYluS_QKYSILRVCJHvY6KpSy2-ds"
genai.configure(api_key=GENMINI_API_KEY)

app = FastAPI()

print("🚀 API_Rag đang khởi động...")

# ✅ Thêm CORS Middleware để React frontend có thể gọi API mà không lỗi CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Hoặc thay "*" bằng ["http://localhost:3000"] nếu chỉ dùng nội bộ
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

print("✅ Cấu hình CORS: Cho phép tất cả nguồn gốc")

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DB_PATH = os.path.join(BASE_DIR, "chroma_db")
MODEL_NAME = "intfloat/multilingual-e5-base"

# Cấu hình đường dẫn uploads
UPLOAD_FOLDER = "uploads"

# Mount thư mục uploads để phục vụ các file tĩnh
app.mount("/static", StaticFiles(directory=UPLOAD_FOLDER), name="static")

# ✅ Kiểm tra model có tải thành công không
try:
    sentence_ef = SentenceTransformer(MODEL_NAME)
except Exception as e:
    print(f"⚠️ Lỗi khi tải model {MODEL_NAME}: {e}")

# ✅ Kiểm tra kết nối ChromaDB
try:
    chroma_client = chromadb.PersistentClient(path=DB_PATH)
    collection = chroma_client.get_or_create_collection(name="my_collection")
except Exception as e:
    print(f"⚠️ Lỗi khi kết nối ChromaDB: {e}")

# Định nghĩa request
class QuestionRequest(BaseModel):
    question: str

class GradeRequest(BaseModel):
    student_image_path: str
    answer_key: str

# Tạo class TopicRequest cho API tạo câu hỏi từ admin
class TopicRequest(BaseModel):
    topic: str

# Định nghĩa model cho dữ liệu export Excel
class ExcelExportRequest(BaseModel):
    results: list

def get_fresh_collection():
    from chromadb import PersistentClient
    from chromadb.config import Settings

    chroma_client = PersistentClient(path=DB_PATH)
    return chroma_client.get_collection(name="my_collection")

def search_similar_chunks(question, top_k=3):
    try:
        print(f"🔍 Searching for: '{question}', top_k={top_k}")
        query_embedding = sentence_ef.encode([question]).tolist()
        
        # Đếm số lượng văn bản trong collection
        collection = get_fresh_collection()  # 🔁 luôn lấy bản cập nhật mới nhất từ disk
        collection_info = collection.get(include=["metadatas", "documents"])
        total_chunks = len(collection_info["ids"]) if "ids" in collection_info else 0
        print(f"💾 Database contains {total_chunks} total chunks")
        
        # In ra tất cả ID để debug
        if "ids" in collection_info and collection_info["ids"]:
            print(f"💾 First 5 IDs in database: {collection_info['ids'][:5]}")
            if "metadatas" in collection_info and collection_info["metadatas"]:
                for i in range(min(len(collection_info["metadatas"]), 5)):
                    print(f"Metadata {i}: {collection_info['metadatas'][i]}")
        
        if total_chunks == 0:
            print("⚠️ No documents in the database. Please add some documents first.")
            return []

        # ===== SEMANTIC SEARCH =====
        # Search với tất cả các chunk có trong database
        search_top_k = total_chunks  
        print(f"🔍 Performing semantic search with top_k={search_top_k}")
        
        # Lấy lại collection mới nhất trước khi thực hiện query để đảm bảo dữ liệu mới nhất  
        try:
            results = collection.query(
                query_embeddings=query_embedding, 
                n_results=search_top_k
            )
            
            found_count = len(results['documents'][0]) if results['documents'] else 0
            print(f"🔍 Semantic search results: found {found_count} documents")
        except Exception as query_error:
            print(f"⚠️ Error during semantic search: {str(query_error)}")
            results = {"documents": [], "metadatas": [], "distances": []}
        
        # ===== KEYWORD SEARCH =====
        # Thêm tìm kiếm đơn giản dựa trên từ khóa
        print(f"🔍 Performing keyword search for terms in: '{question}'")
        keywords = [kw.lower() for kw in question.lower().split() if len(kw) > 2]  # Ignore very short words
        print(f"🔍 Keywords: {keywords}")
        keyword_results = []
        
        # Duyệt qua tất cả các documents từ collection_info (đã lấy ở trên)
        if "documents" in collection_info and collection_info["documents"]:
            for i, (doc_id, doc, metadata) in enumerate(zip(
                collection_info["ids"], 
                collection_info["documents"], 
                collection_info["metadatas"]
            )):
                if metadata is None:
                    print(f"⚠️ Skipping document with None metadata (ID: {doc_id})")
                    continue
                    
                # Tính điểm match đơn giản: số từ khóa được tìm thấy trong document
                doc_lower = doc.lower()
                keyword_matches = sum(1 for kw in keywords if kw in doc_lower)
                
                # Kiểm tra title nếu có
                title_matches = 0
                if "title" in metadata and metadata["title"]:
                    title_lower = metadata["title"].lower() 
                    title_matches = sum(1 for kw in keywords if kw in title_lower)
                
                total_matches = keyword_matches + (title_matches * 2)  # Title matches count more
                
                if total_matches > 0:
                    # Nếu có ít nhất một từ khóa match, thêm vào kết quả
                    keyword_score = 1.0 - (total_matches / (len(keywords) * 3))  # Điểm thấp = match tốt
                    keyword_results.append({
                        'id': doc_id,
                        'document': doc,
                        'metadata': metadata,
                        'score': keyword_score,
                        'keyword_matches': total_matches
                    })
                    print(f"📄 Keyword match: {metadata.get('title', 'Unnamed Chunk')} (matches: {total_matches})")
        
        # Sắp xếp keyword results theo số lượng keywords match (giảm dần)
        keyword_results.sort(key=lambda x: x['score'])
        print(f"🔍 Keyword search results: found {len(keyword_results)} documents")
        
        # ===== COMBINE RESULTS =====
        # Tạo một dictionary để track chunks đã được chọn bởi ID
        combined_chunks = {}
        
        # Thêm semantic search results
        if results["documents"] and len(results["documents"]) > 0 and len(results["documents"][0]) > 0:
            try:
                # Đảm bảo metadatas và documents có cùng độ dài
                doc_count = len(results["documents"][0])
                meta_count = len(results["metadatas"][0]) if results["metadatas"] else 0
                dist_count = len(results["distances"][0]) if results["distances"] else 0
                
                print(f"🔍 Debug: documents: {doc_count}, metadatas: {meta_count}, distances: {dist_count}")
                
                # Chỉ xử lý các phần tử có đủ thông tin
                for i in range(doc_count):
                    if i >= meta_count or i >= dist_count:
                        print(f"⚠️ Index {i} exceeds available metadata or distances")
                        continue
                    
                    doc = results["documents"][0][i]
                    metadata = results["metadatas"][0][i] if results["metadatas"] else None
                    score = results["distances"][0][i] if results["distances"] else 1.0
                    
                    if metadata is None:
                        print(f"⚠️ Skipping result with None metadata (index {i})")
                        # Tạo metadata tạm thời dựa trên nội dung doc
                        title = doc[:50].replace("\n", " ")
                        if len(title) >= 50:
                            title += "..."
                        
                        # Tạo metadata mới với ID dựa trên hash của nội dung
                        import hashlib
                        doc_hash = hashlib.md5(doc.encode()).hexdigest()[:10]
                        temp_doc_id = f"temp_{doc_hash}"
                        
                        metadata = {
                            'page': 0, 
                            'document_id': temp_doc_id,
                            'title': title,
                            'chunk_id': f"{temp_doc_id}_chunk_0"
                        }
                        print(f"🔄 Created temporary metadata for document: {title}")
                    
                    chunk_id = metadata.get('chunk_id', '')
                    if not chunk_id:
                        print(f"⚠️ Missing chunk_id in metadata")
                        continue
                        
                    if chunk_id not in combined_chunks:
                        combined_chunks[chunk_id] = {
                'page': metadata.get('page', metadata.get('chunk_index', 0)), 
                'content': doc, 
                'score': score,
                'document_id': metadata.get('document_id', 'unknown'),
                            'title': metadata.get('title', 'Unnamed Chunk'),
                            'source': 'semantic'
                        }
                        print(f"📄 Added from semantic: {metadata.get('title', 'Unnamed Chunk')} (score: {score:.4f})")
                    
            except Exception as combine_error:
                print(f"⚠️ Error combining semantic results: {str(combine_error)}")
                import traceback
                traceback.print_exc()
        
        # Thêm keyword search results
        for result in keyword_results:
            try:
                metadata = result['metadata']
                if metadata is None:
                    print(f"⚠️ Skipping keyword result with None metadata")
                    continue
                    
                chunk_id = metadata.get('chunk_id', '')
                if not chunk_id:
                    print(f"⚠️ Missing chunk_id in keyword result metadata")
                    continue
                    
                if chunk_id not in combined_chunks:
                    combined_chunks[chunk_id] = {
                        'page': metadata.get('page', metadata.get('chunk_index', 0)), 
                        'content': result['document'], 
                        'score': result['score'],  # Sử dụng score từ keyword search
                        'document_id': metadata.get('document_id', 'unknown'),
                        'title': metadata.get('title', 'Unnamed Chunk'),
                        'source': 'keyword',
                        'keyword_matches': result['keyword_matches']
                    }
                    print(f"📄 Added from keyword: {metadata.get('title', 'Unnamed Chunk')} (score: {result['score']:.4f}, matches: {result['keyword_matches']})")
            except Exception as kw_error:
                print(f"⚠️ Error adding keyword result: {str(kw_error)}")
        
        # Chuyển đổi từ dictionary thành list
        chunks_found = list(combined_chunks.values())
        
        # Sắp xếp kết quả: semantic trước, sau đó đến keyword, rồi theo score
        chunks_found.sort(key=lambda x: (
            0 if x.get('source') == 'semantic' else 1,  # semantic trước
            x.get('score', 1.0)  # score thấp hơn trước (tốt hơn)
        ))
        
        print(f"✅ Combined results: {len(chunks_found)} unique chunks")
        
        # Print details of combined results
        for i, chunk in enumerate(chunks_found):
            print(f"📄 Result {i+1}: {chunk['title']} (score: {chunk['score']:.4f}, source: {chunk['source']})")
        
        # Nếu không tìm thấy kết quả nào
        if not chunks_found:
            print("⚠️ No chunks found after combining results")
            # Tìm kiếm đơn giản hơn - chỉ dựa trên từ khóa chính
            main_keywords = [kw for kw in keywords if len(kw) > 3]
            if main_keywords:
                main_keyword = main_keywords[0]
                print(f"🔎 Performing fallback search with main keyword: '{main_keyword}'")
                
                for doc_id, doc, metadata in zip(
                    collection_info["ids"], 
                    collection_info["documents"],
                    collection_info["metadatas"]
                ):
                    if metadata is None:
                        continue
                        
                    if main_keyword in doc.lower():
                        print(f"🔎 Found document containing '{main_keyword}'")
                        chunks_found.append({
                            'page': metadata.get('page', metadata.get('chunk_index', 0)), 
                            'content': doc, 
                            'score': 0.5,  # Score trung bình
                            'document_id': metadata.get('document_id', 'unknown'),
                            'title': metadata.get('title', 'Unnamed Chunk'),
                            'source': 'fallback'
                        })
                        break
                        
        # Trả về top_k kết quả hoặc tất cả những gì tìm được nếu ít hơn top_k
        final_results = chunks_found[:top_k] if len(chunks_found) > 0 else []
        print(f"✅ Returning {len(final_results)} chunks")
        return final_results
    except Exception as e:
        print(f"⚠️ Lỗi khi tìm kiếm tài liệu: {str(e)}")
        import traceback
        traceback.print_exc()
        return []

def generate_answer(question, context):
    try:
        model_gen = genai.GenerativeModel('gemini-2.0-flash')  # 🔄 Cập nhật model mới hơn nếu có
        input_parts = [
            f"""
            Bạn là một trợ lý AI gia sư Toán học với nhiều năm kinh nghiệm giảng dạy.  
            Nhiệm vụ của bạn là hướng dẫn học sinh một cách tự nhiên, rõ ràng và dễ hiểu, **chỉ dựa trên thông tin từ sách giáo khoa mà tôi cung cấp**.  

            ---  

            ### **🚨 Lưu ý quan trọng:**  
            - **KHÔNG** sử dụng bất kỳ kiến thức nào từ nguồn ngoài – chỉ dựa trên nội dung sách giáo khoa. 
            - ** Nếu câu hỏi không liên quan đến chủ đề Toán học thì từ chối trả lờilời"** 
            - **Nếu câu hỏi liên quan đến toán học nhưng kiến thức trong sách giao khoa không có đủ thông tin liên quan để trả lời câu hỏi thì trả lời là ""Không đủ kiến thứuc để trả lời"
            - **Trả về kết quả định dạng Markdown**, đảm bảo hiển thị tốt trong React:
            - **Công thức toán inline:** `$...$`
            - **Công thức toán block:** `$$...$$`  
            - **Chữ in đậm:** `**...**`

            ---  

            ### **📚 Thông tin từ sách giáo khoa:**  
            {context}  

            ### **❓ Câu hỏi của học sinh:**  
            {question}  

            ### **🔹 Hướng dẫn trả lời (tuân theo cấu trúc sau):**  

            #### **1. Tóm tắt kiến thức trọng tâm:**  
            - Trình bày ngắn gọn các khái niệm, công thức và phương pháp giải.
            - Các công thức phải được viết dưới dạng LaTeX:
            - **Inline:** `Diện tích hình tròn: $S = \pi r^2$`
            - **Block:**  
                ```md
                $$S = \pi r^2$$
                ```

            #### **2. Bản tóm tắt và lưu ý quan trọng:**  
            - Liệt kê các điểm mấu chốt, lỗi sai thường gặp.

            #### **3. Các thuật ngữ quan trọng:**  
            - Giải thích các thuật ngữ chuyên môn bằng ngôn ngữ dễ hiểu.

            #### **. Trả lời câu hỏi:**  
            - Trả lời câu hỏi trực tiếp dựa trên kiến thức đã trình bày.

            ---

            ### **⚠️ Nguyên tắc quan trọng khi trả lời:**  

            ✅ **Chỉ dựa trên thông tin từ sách giáo khoa.**  
            - Nếu sách không cung cấp đủ thông tin, trả lời:  
            > _"Tôi không có đủ thông tin để trả lời câu hỏi này dựa trên sách giáo khoa."_  
            -

            ✅ **Phong cách giảng dạy của gia sư:**  
            - Sử dụng lời văn tự nhiên, gần gũi.
            - Dùng phép so sánh trực quan để giúp học sinh dễ hình dung.
            - Đặt câu hỏi ngược lại để kích thích tư duy học sinh.

            ✅ **Hướng dẫn từng bước giải bài tập:**  
            - KHÔNG giải bài tập ngay lập tức; hướng dẫn từng bước để học sinh tự suy nghĩ.
            - Đưa ra bài tập bổ sung có độ khó tăng dần.

            ---

            📈 **Hãy đảm bảo câu trả lời của bạn rõ ràng, có cấu trúc khoa học(xuống dòng hợp lý) và giúp học sinh phát triển tư duy độc lập!**

        """
        ]

        response = model_gen.generate_content(input_parts)
        return response.text.strip() if response else "❌ Không có phản hồi từ GenMini."
    except Exception as e:
        print(f"⚠️ Lỗi khi gọi GenMini API: {e}")
        return "⚠️ Lỗi khi gọi GenMini API."
    
# //-----------------------------------------------------------------------------------------------

def generate_Multiple_Choice_Questions(question, context):
    try:
        model_gen = genai.GenerativeModel('gemini-2.0-flash')  # 🔄 Cập nhật model mới hơn nếu có
        input_parts = [
            f"""
            # Hướng Dẫn Tạo Câu Hỏi Trắc Nghiệm Toán Học Thông Minh với Đánh Giá Độ Khó

## 🤖 Vai Trò
Bạn là một trợ lý AI chuyên tạo câu hỏi trắc nghiệm Toán học và đánh giá độ khó của các câu hỏi.

## 📚 Đầu Vào
📖 Thông tin từ sách giáo khoa:
""" + context + """

❓ Câu hỏi của học sinh:
""" + question + """

## 🎯 Nguyên Tắc Tạo Câu Hỏi
1. Yêu Cầu Cơ Bản:
   - Liên quan trực tiếp đến chủ đề Câu hỏi của học sinh.
   - Xây dựng dựa trên thông tin từ sách giáo khoa.
   - Có độ khó đa dạng phân bố ở cả 3 mức độ khó

2. Chiến Lược Chi Tiết
   - Phân tích sâu câu hỏi của học sinh và thông tin từ sách giáo khoa.
   - Xác định từ khóa chính
   - Tạo câu hỏi từ nhận biết đến vận dụng cao

3. **Lưu ý quan trọng:**  
   - **Trả về kết quả định dạng Markdown**, đảm bảo hiển thị tốt trong React:
   - **Công thức toán inline:** `$...$`
   - **Công thức toán block:** `$$...$$`  
   - **Chữ in đậm:** `**...**`
   - Các công thức phải được viết dưới dạng LaTeX:
   - **Inline:** `Diện tích hình tròn: $S = \pi r^2$`
   - **Block:**  
     ```md
     $$S = \pi r^2$$
     ```
   - Không thêm khoảng trắng thừa
   - Sử dụng các lệnh LaTeX chuẩn

## 📝 Cấu Trúc Câu Hỏi
- Tổng: **10 câu hỏi**
- Mỗi câu: 4 đáp án (A, B, C, D)
- Chỉ 1 đáp án đúng
- Phân bố: 3-4 câu Mức 1, 3-4 câu Mức 2, 2-3 câu Mức 3

## 🔢 Tiêu Chí Đánh Giá Độ Khó
### Mức 1 - Cơ bản
- Yêu cầu áp dụng trực tiếp công thức, định nghĩa hoặc quy tắc cơ bản
- Chỉ cần một bước giải đơn giản hoặc vài bước đơn giản
- Không đòi hỏi biến đổi phức tạp
- Sử dụng kiến thức cơ bản trong chương trình
- Học sinh trung bình có thể giải quyết trong thời gian ngắn

### Mức 2 - Trung bình
- Yêu cầu kết hợp 2-3 công thức hoặc khái niệm
- Cần nhiều bước giải quyết có logic
- Có thể có một số biến đổi toán học vừa phải
- Đòi hỏi hiểu sâu về kiến thức trong chương trình
- Học sinh khá có thể giải quyết được

### Mức 3 - Nâng cao
- Yêu cầu kết hợp nhiều công thức, khái niệm từ các phần khác nhau
- Đòi hỏi nhiều bước giải với cách tiếp cận sáng tạo
- Có các biến đổi toán học phức tạp
- Cần tư duy phân tích, tổng hợp hoặc suy luận logic cao
- Chỉ học sinh giỏi mới có thể giải quyết được
- Có thể chứa nội dung mở rộng hoặc nâng cao

## 🔍 Định Dạng JSON
```json
{
  "questions": [
    {
      "question": "Nội dung câu hỏi",
      "options": {
        "A": "Đáp án A",
        "B": "Đáp án B", 
        "C": "Đáp án C",
        "D": "Đáp án D"
      },
      "answer": "B",
      "difficulty": "Mức 1/Mức 2/Mức 3",
      "solution": "Lời giải chi tiết từng bước, giải thích công thức và phân tích sai lầm thường gặp nếu có "
    }
            ]
            }
            ```
    ⚠️ Lưu Ý QUAN TRỌNG

    Nếu không đủ thông tin: Trả về JSON rỗng
    Ưu tiên thuật ngữ từ text trong sách giáo khoa
    LUÔN LUÔN tuân thủ quy tắc viết công thức toán học
    TRÁNH sinh những câu hỏi cần nhìn hình hoặc dựa vào bảng biến thiên
    ĐÁP ÁN VÀ LỜI GIẢI PHẢI ĐÚNG VÀ ĂN KHỚP VỚI NHAU.CHỈ ĐƯỢC PHÉP CÓ 1 ĐÁP ÁN ĐÚNG
    Tránh lỗi sai định dạng JSON
    Lưu ý thêm xuống dòng phù hợp để trình bày đẹp hơn
    Phải sinh đủ 10 câu hỏi không được thiếu câu nào.
    """
        ]

        response = model_gen.generate_content(input_parts)
        return response.text.strip() if response else "❌ Không có phản hồi từ GenMini."
    except Exception as e:
        print(f"⚠️ Lỗi khi gọi GenMini API: {e}")
        return "⚠️ Lỗi khi gọi GenMini API."

# //-----------------------------------------------------------------------------------------------
# sinh câu hỏi lý thuyết(Chưa sử dụng)
def generate_Essay(question, context):
    try:
        model_gen = genai.GenerativeModel('gemini-2.0-flash')  # 🔄 Cập nhật model mới hơn nếu có
        input_parts = [
            f"""
            Bạn là một trợ lý AI chuyên tạo câu hỏi lý thuyết Toán học.

            📚 Đầu Vào
            📖 Thông tin từ sách giáo khoa:
            """ + context + """

            ❓ Câu hỏi của học sinh:
            """ + question + """

            🎯 Nguyên Tắc Tạo Câu Hỏi
            1. Yêu Cầu Cơ Bản:
            Câu hỏi ngắn gọn, tập trung vào lý thuyết (khái niệm, định nghĩa, công thức, tính chất).

            Không yêu cầu tính toán hay chọn đáp án.

            Có sẵn câu trả lời, không cần người học tự suy luận hoặc giải bài toán.

            2. Chiến Lược Chi Tiết
            Chọn lọc thông tin quan trọng từ sách giáo khoa.

            Tạo câu hỏi đơn giản, dễ hiểu, có câu trả lời chính xác.

            Tránh câu hỏi quá rộng hoặc cần minh họa bằng hình vẽ.

            3. Lưu ý quan trọng:
            Trả về kết quả định dạng Markdown, đảm bảo hiển thị tốt trong React.

            Công thức toán inline: $...$

            Công thức toán block: $$...$$

            Chữ in đậm: **...**

            Các công thức phải được viết dưới dạng LaTeX.

            📝 Cấu Trúc Câu Hỏi
            Tổng: 10 câu hỏi

            Mỗi câu hỏi có dạng:

            Câu hỏi: Nội dung câu hỏi

            Đáp án: Câu trả lời chính xác

            🔍 Định Dạng JSON
            json

            {
            "questions": [
                {
                "question": "Nội dung câu hỏi lý thuyết",
                "answer": "Câu trả lời chính xác"
                }
            ]
            }
            ⚠️ Lưu Ý QUAN TRỌNG
            Nếu không đủ thông tin: Trả về JSON rỗng

            Chỉ tạo câu hỏi lý thuyết, không yêu cầu tính toán hay trắc nghiệm.

            Không sinh câu hỏi cần hình minh họa hoặc bảng số liệu.

            LUÔN tuân thủ định dạng JSON chuẩn, không để lỗi cú pháp.
    """
        ]

        response = model_gen.generate_content(input_parts)
        return response.text.strip() if response else "❌ Không có phản hồi từ GenMini."
    except Exception as e:
        print(f"⚠️ Lỗi khi gọi GenMini API: {e}")
        return "⚠️ Lỗi khi gọi GenMini API."
    
# //--------------------------------------------------------------------------------
# hàm này tìm topic từ các lượt trò chuyện 
def get_chat_topic(context):
    try:
        model_gen = genai.GenerativeModel('gemini-2.0-flash')  # 🔄 Cập nhật model mới hơn nếu có
        input_parts = [
            f"""
            Đây là các lượt hỏi gần nhất của người dùng cho hệ thống: 

            {context}

            Dựa trên các lượt hoir, hãy xác định chủ đề chính mà người dùng đang quan tâm. 
            Chỉ cần đưa ra một chủ đề duy nhất mà người dùng đang thảo luận.
        """
        ]

        response = model_gen.generate_content(input_parts)
        return response.text.strip() if response else "❌ Không có phản hồi từ GenMini."
    except Exception as e:
        print(f"⚠️ Lỗi khi gọi GenMini API: {e}")
        return "⚠️ Lỗi khi gọi GenMini API."

# Hàm này trả về feedback về bài kiểm tra của sinh viên
def get_Quiz_Feedback(context):
    try:
        model_gen = genai.GenerativeModel('gemini-2.0-flash') 
        input_parts = [
            f"""
           Bạn là một gia sư toán học. Tôi sẽ cung cấp dữ liệu bài kiểm tra của một học sinh dưới dạng JSON (chứa câu hỏi, đáp án đúng, câu trả lời của học sinh và điểm số).
            Hãy phân tích bài làm của học sinh, nhận xét về điểm mạnh, điểm yếu của học sinh và đề xuất hướng cải thiện.
            Trả lời theo định dạng sau:

            Điểm số: [ghi điểm số của học sinh]

            Nhận xét tổng quan: [đưa ra đánh giá tổng quan về bài làm]

            Điểm mạnh: [liệt kê những điểm học sinh làm tốt]

            Điểm yếu: [liệt kê những lỗi sai, kiến thức cần cải thiện]

            Gợi ý ôn tập: [đề xuất nội dung hoặc dạng bài tập học sinh nên luyện tập thêm]

            Dữ liệu bài làm của học sinh:
            {context}
            với độ khó được đánh giá như sau:
            ### Mức 1 - Cơ bản
            - Yêu cầu áp dụng trực tiếp công thức, định nghĩa hoặc quy tắc cơ bản
            - Chỉ cần một bước giải đơn giản hoặc vài bước đơn giản
            - Không đòi hỏi biến đổi phức tạp
            - Sử dụng kiến thức cơ bản trong chương trình
            - Học sinh trung bình có thể giải quyết trong thời gian ngắn

            ### Mức 2 - Trung bình
            - Yêu cầu kết hợp 2-3 công thức hoặc khái niệm
            - Cần nhiều bước giải quyết có logic
            - Có thể có một số biến đổi toán học vừa phải
            - Đòi hỏi hiểu sâu về kiến thức trong chương trình
            - Học sinh khá có thể giải quyết được

            ### Mức 3 - Nâng cao
            - Yêu cầu kết hợp nhiều công thức, khái niệm từ các phần khác nhau
            - Đòi hỏi nhiều bước giải với cách tiếp cận sáng tạo
            - Có các biến đổi toán học phức tạp
            - Cần tư duy phân tích, tổng hợp hoặc suy luận logic cao
            - Chỉ học sinh giỏi mới có thể giải quyết được
            - Có thể chứa nội dung mở rộng hoặc nâng cao
        """
        ]

        response = model_gen.generate_content(input_parts)
        return response.text.strip() if response else "❌ Không có phản hồi từ GenMini."
    except Exception as e:
        print(f"⚠️ Lỗi khi gọi GenMini API: {e}")
        return "⚠️ Lỗi khi gọi GenMini API."
# Hàm này gọi khi sinh lại câu hỏi( ấn làm lại sau khi làm bài trắc nghiệm theo chủ đềđề)
def Adaptive_Questions(topic, previous_results_json):
    try:
        model_gen = genai.GenerativeModel('gemini-2.0-flash') 
        input_parts = [
            f"""
           # Hướng Dẫn Tạo Câu Hỏi Trắc Nghiệm Toán Học Adaptive Dựa Trên Bài Làm Trước Đó

            ## 🤖 Vai Trò
            Bạn là một trợ lý AI chuyên tạo câu hỏi trắc nghiệm Toán học thông minh (adaptive), giúp học sinh luyện tập hiệu quả hơn dựa vào kết quả làm bài trước đó.

            ## 📚 Đầu Vào
            📖 Chủ đề trọng tâm:  
            """ + topic + """

            🧾 Thông tin bài làm trước đó của học sinh:
            ```json
            """ + previous_results_json + """
            ```

            ## 🎯 Mục Tiêu
            1. **Củng cố các phần học sinh yếu** (học sinh câu sai, đặc biệt nếu sai nhiều lần hoặc sai ở Mức 1-2). thì giảm độ khó của đề bằng cách tăng câu hỏi mức 1-2 lên
            2. **Tăng thử thách cho phần học sinh làm tốt** (làm đúng nhiều lần, đúng nhanh và đúng ở mức độ 3 nhiều) thì tăng số câu hỏi mức 3 lên
            3. **Câu hỏi chỉ nên cùng dạng không được giống ý hệt câu hỏi mà học sinh đã làm trước đó.**
            4. **ĐÁP ÁN VÀ LỜI GIẢI PHẢI ĐÚNG VÀ ĂN KHỚP VỚI NHAU.CHỈ ĐƯỢC PHÉP CÓ 1 ĐÁP ÁN ĐÚNG**
            ## 🔧 Nguyên Tắc Tạo Câu Hỏi
            - Mỗi câu hỏi phải liên quan trực tiếp đến **chủ đề trọng tâm** và xoay quanh những câu sai nhiều hoặc phải khó hơn nếu học sinh làm đúngđúng.
            - Có 4 đáp án A, B, C, D, chỉ một đáp án đúng.
            - Mỗi câu phải có trường `"difficulty"` là `"Mức 1"`, `"Mức 2"` hoặc `"Mức 3"`.
            - Chú ý xuống dòng phù hợp để trình bày đẹp hơn
            - **Tránh sinh câu hỏi yêu cầu nhìn hình, bảng biến thiên hoặc đồ thị.**
            5. Lưu ý quan trọng:
                    Trả về kết quả định dạng Markdown, đảm bảo hiển thị tốt trong React.

                    Công thức toán inline: $...$

                    Công thức toán block: $$...$$

                    Chữ in đậm: **...**

                    Các công thức phải được viết dưới dạng LaTeX.

                    📝 Cấu Trúc Câu Hỏi
                    Tổng: 10 câu hỏi

                    Mỗi câu hỏi có dạng:

                    Câu hỏi: Nội dung câu hỏi

                    Đáp án: Câu trả lời chính xác

            ## 📝 Định Dạng Kết Quả
            ```json
            {
            "questions": [
                {
                "question": "Nội dung câu hỏi",
                "options": {
                    "A": "Đáp án A",
                    "B": "Đáp án B", 
                    "C": "Đáp án C",
                    "D": "Đáp án D"
                },
                "answer": "B",
                "difficulty": "Mức 1/Mức 2/Mức 3",
                "solution": "Lời giải chi tiết từng bước, giải thích công thức và phân tích sai lầm thường gặp nếu có"
                }
            ]
            }
            ```

            ## ✅ Lưu Ý
            - Các công thức toán học phải viết bằng LaTeX:
            - Inline: `$...$`
            - Block:  
                ```md
                $$...$$
                ```
            - Kết quả trả về phải đúng chuẩn JSON.
            - Nếu thiếu thông tin: Trả về `{ "questions": [] }`

        """
        ]

        response = model_gen.generate_content(input_parts)
        return response.text.strip() if response else "❌ Không có phản hồi từ GenMini."
    except Exception as e:
        print(f"⚠️ Lỗi khi gọi GenMini API: {e}")
        return "⚠️ Lỗi khi gọi GenMini API."

def Grade_math_paper(student_image_path, answer_key, model_name='gemini-2.0-flash'):
    
    # Function to encode image to base64
    def encode_image(image_path):
        try:
            # Chuẩn hóa đường dẫn
            image_path = os.path.normpath(image_path)
            print(f"Đang đọc file: {image_path}")
            
            # In thông tin về đường dẫn để debug
            abs_path = os.path.abspath(image_path)
            print(f"Đường dẫn tuyệt đối: {abs_path}")
            
            # Kiểm tra thư mục chứa file có tồn tại không
            parent_dir = os.path.dirname(image_path)
            if not os.path.exists(parent_dir):
                print(f"Thư mục cha không tồn tại: {parent_dir}")
                return None
            
            # Kiểm tra file có tồn tại không
            if not os.path.exists(image_path):
                print(f"File không tồn tại: {image_path}")
                # Đợi một chút để đảm bảo file system đã cập nhật
                time.sleep(2)
                # Kiểm tra lại
                if not os.path.exists(image_path):
                    print(f"File vẫn không tồn tại sau khi đợi: {image_path}")
                    return None
                print(f"File đã tồn tại sau khi đợi: {image_path}")
            
            # Kiểm tra kích thước file
            file_size = os.path.getsize(image_path)
            print(f"Kích thước file: {file_size} bytes")
            
            if file_size == 0:
                print(f"File rỗng (0 bytes): {image_path}")
                time.sleep(2)  # Đợi lâu hơn để file có thể được ghi đầy đủ
                file_size = os.path.getsize(image_path)
                if file_size == 0:
                    print(f"File vẫn rỗng sau khi đợi: {image_path}")
                    return None
                print(f"File đã có dữ liệu sau khi đợi: {image_path}, kích thước: {file_size} bytes")
            
            # Đọc file và mã hóa base64
            with open(image_path, "rb") as image_file:
                file_content = image_file.read()
                if not file_content:
                    print(f"Đọc được file nhưng nội dung rỗng: {image_path}")
                    return None
                encoded = base64.b64encode(file_content).decode('utf-8')
                print(f"Đã mã hóa file thành công: {image_path} ({len(encoded)} ký tự)")
                return encoded
                
        except Exception as e:
            print(f"Lỗi khi xử lý file: {e}")
            return None
    
    # Process student image
    student_image_b64 = encode_image(student_image_path)
    if not student_image_b64:
        return f"❌ Failed to process student's paper image. Path: {student_image_path}"
    
    # Process answer key - check if it's an image path or text
    is_answer_key_image = os.path.isfile(answer_key) if isinstance(answer_key, str) else False
    answer_key_content = ""
    answer_key_b64 = None
    
    if is_answer_key_image:
        answer_key_b64 = encode_image(answer_key)
        if not answer_key_b64:
            return f"❌ Failed to process answer key image. Path: {answer_key}"
    else:
        answer_key_content = answer_key
    
    try:
        model_gen = genai.GenerativeModel(model_name)  # Load the model
        
        # Prepare the prompt
        prompt = """
        # Hướng dẫn chấm điểm bài làm toán viết tay của học sinh

        ## 🧑‍🏫 Vai trò của bạn
        Bạn là một giáo viên toán giàu kinh nghiệm, chấm điểm bài làm viết tay của học sinh theo biểu điểm cho sẵn một cách công bằng, chính xác và chi tiết.

        ## 📝 Nhiệm vụ
        1. Nếu ảnh bài làm của sinh viên không liên quan đến bài làm thì chỉ chả về duy nhất là bài làm lạc đề
        2. Phân tích kỹ hình ảnh bài làm viết tay của học sinh
        3. So sánh với đáp án và biểu điểm được cung cấp
        4. Chấm điểm chi tiết từng câu, từng ý
        5. Nêu rõ lỗi sai và thiếu sót (nếu có)
        6. Đề xuất hướng cải thiện

        ## 🎯 Yêu cầu kết quả
        1. **Điểm số cụ thể** cho từng câu và tổng điểm
        2. **Nhận xét chi tiết** về từng câu:
           - Đúng hoàn toàn: nêu rõ các bước làm đúng
           - Đúng một phần: nêu rõ điểm đúng và điểm thiếu
           - Sai: phân tích lỗi sai và cách sửa
        3. **Nhận xét tổng quát** về bài làm, ưu điểm và nhược điểm
        4. **Đề xuất cụ thể** để học sinh tiến bộ

        ## 📋 Trình bày kết quả(Lưu ý chỉ trình bày kết quả nếu như bài làm của sinh viên liên quan đề đề bài.Nếu không thì bỏ qua phần này)
        Trả lời theo cấu trúc sau:
        ```
        # KẾT QUẢ CHẤM ĐIỂM
        ## THÔNG TIN SINH VIÊN
        -Họ và tên: Lê văn x
        -Lớp: 12A
        ## ĐIỂM SỐ
        - Câu 1: x/y điểm
        - Câu 2: x/y điểm
        ...
        - TỔNG ĐIỂM: X/Y

        ## NHẬN XÉT CHI TIẾT
        ### Câu 1:
        - Nhận xét...
        - Lỗi sai/thiếu sót (nếu có)...

        ### Câu 2:
        - Nhận xét...
        - Lỗi sai/thiếu sót (nếu có)...
        ...

        ## NHẬN XÉT TỔNG QUÁT
        - Ưu điểm:...
        - Nhược điểm:...

        ## ĐỀ XUẤT CẢI THIỆN
        1. ...
        2. ...
        ```

        ## ⚠️ Lưu ý quan trọng
        - Phải có đầy đủ các trường thông tin cần thiết( đặc biệt là: -Họ và tên:;-Lớp: 12A; - TỔNG ĐIỂM:)
        - Chấm điểm công bằng, không quá nghiêm khắc hay quá dễ dãi
        - Nhận diện các phương pháp giải khác với đáp án nhưng vẫn đúng
        - Khi học sinh làm đúng kết quả nhưng cách giải khác, cần phân tích cách giải của học sinh xem có chặt chẽ không
        - Với những câu khó hiểu hoặc không rõ ràng, hãy nêu rõ trong phần nhận xét
        - Đề xuất cải thiện phải cụ thể, dựa trên lỗi sai thường gặp của học sinh
        """
        
        # Add content based on the type of answer key
        parts = [{"text": prompt}]
        
        # Add student's paper image
        parts.append({
            "inline_data": {
                "mime_type": "image/jpeg",
                "data": student_image_b64
            }
        })
        
        # Add answer key information
        if is_answer_key_image:
            parts.append({
                "text": "\n\n## 📑 Đáp án và biểu điểm (trong hình ảnh):"
            })
            parts.append({
                "inline_data": {
                    "mime_type": "image/jpeg",
                    "data": answer_key_b64
                }
            })
        else:
            parts.append({
                "text": f"\n\n## 📑 Đáp án và biểu điểm:\n\n{answer_key_content}"
            })
        
        # Generate response
        response = model_gen.generate_content(parts)
        return response.text.strip() if response else "❌ Không có phản hồi từ LLM."
    
    except Exception as e:
        print(f"⚠️ Lỗi khi gọi LLM API: {e}")
        return f"⚠️ Lỗi khi chấm điểm: {e}"

def generate_Slide(question, context):
    try:
        model_gen = genai.GenerativeModel('gemini-2.0-flash')
        input_parts = [
            f"""
            # Hướng Dẫn Tạo Slide Bài Giảng Toán Học

## 🤖 Vai Trò
Bạn là một trợ lý AI chuyên tạo nội dung slide bài giảng toán học, giúp giáo viên chuẩn bị giáo án.

## 📚 Đầu Vào
📖 Thông tin từ sách giáo khoa:
""" + context + """

❓ Chủ đề bài giảng:
""" + question + """

## 🎯 Nguyên Tắc Tạo Slide
1. Yêu Cầu Cơ Bản:
   - Tạo 5-10 slide liên quan trực tiếp đến chủ đề
   - Nội dung dựa trên thông tin từ sách giáo khoa
   - Đảm bảo tính sư phạm và logic của bài giảng

2. Chiến Lược Chi Tiết:
   - Mỗi slide có cấu trúc rõ ràng với tiêu đề và nội dung
   - Các slide được sắp xếp theo trình tự logic: giới thiệu khái niệm, phát triển, ví dụ, bài tập
   - Phù hợp với cấp độ học sinh THCS hoặc THPT

3. **Lưu ý quan trọng:**  
   - **Trả về kết quả định dạng JSON**
   - Nội dung slide có thể sử dụng cú pháp markdown, LaTeX
   - Phần LaTeX sẽ được xử lý trong PowerPoint

## 📝 Cấu Trúc JSON
```json
[
  {
    "slide_number": 1,
    "title": "Tiêu đề slide",
    "sections": [
      {
        "heading": "Tiêu đề phần",
        "content": "Nội dung phần"
      },
      ...
    ],
    "notes": "Ghi chú dành cho giáo viên (không bắt buộc)"
  },
  ...
]
```

## 🔢 Các Loại Slide Cần Tạo
1. **Slide Trang Bìa (Slide 1)**:
   - Tiêu đề chủ đề
   - Tên môn học: Toán học
   - Cấp độ phù hợp

2. **Slide Mục Tiêu Bài Học**:
   - Kiến thức học sinh sẽ đạt được
   - Kỹ năng sẽ rèn luyện

3. **Slide Nội Dung**:
   - Khái niệm, định nghĩa
   - Công thức, tính chất
   - Ví dụ minh họa
   - Cách giải quyết bài toán

4. **Slide Ví Dụ**:
   - Ví dụ từ đơn giản đến phức tạp
   - Phương pháp giải chi tiết

5. **Slide Bài Tập**:
   - Bài tập áp dụng
   - Bài tập nâng cao (tùy chọn)

6. **Slide Tổng Kết**:
   - Ôn lại kiến thức chính
   - Kết nối với bài học tiếp theo

⚠️ Lưu Ý QUAN TRỌNG:
- Nếu không đủ thông tin: Trả về thông báo lỗi
- Ưu tiên sử dụng thuật ngữ từ sách giáo khoa
- Đảm bảo cấu trúc JSON chính xác
- Nội dung phải khoa học, chính xác về mặt toán học
- các công thức nếu có phù hợp để hiện thị trong powpoin. tránh sinh code dạng ( ký hiệu LaTeX như \(, \))
            """
        ]

        response = model_gen.generate_content(input_parts)
        return response.text.strip() if response else "❌ Không có phản hồi từ GenMini."
    except Exception as e:
        print(f"⚠️ Lỗi khi gọi GenMini API: {e}")
        return "⚠️ Lỗi khi gọi GenMini API."

# API gọi tới hàm generate_answer để trả lời câu hỏihỏi
@app.post("/answer")
async def question(request: QuestionRequest):
    question = request.question
    print(f"💬 Câu hỏi từ người dùng: '{question}'")
    print(f"💻 Request từ IP: {request.client.host if hasattr(request, 'client') else 'Unknown'}")
    
    # In thông tin về collection
    try:
        collection_info = collection.get(include=["metadatas"])
        total_chunks = len(collection_info["ids"]) if "ids" in collection_info else 0
        print(f"💾 Collection hiện có {total_chunks} chunks")
        
        if total_chunks > 0:
            print(f"💾 ID đầu tiên: {collection_info['ids'][0]}")
            print(f"💾 Metadata đầu tiên: {collection_info['metadatas'][0]}")
    except Exception as e:
        print(f"❌ Lỗi khi lấy thông tin từ collection: {str(e)}")
    
    # Thử tìm kiếm với threshold thấp hơn
    retrieved_chunks = search_similar_chunks(question, top_k=5) # Tăng số lượng kết quả

    if not retrieved_chunks:
        print("❌ Không tìm thấy chunks phù hợp trong database")
        return {"question": question, "answer": "❌ Không tìm thấy tài liệu phù hợp. Vui lòng thử lại sau khi tài liệu đã được embed."}

    print(f"✅ Tìm thấy {len(retrieved_chunks)} chunks phù hợp")
    context = "\n\n".join(f"{chunk['content']}" for chunk in retrieved_chunks)
    answer = generate_answer(question, context)
    print(f"✅ Đã tạo câu trả lời dài {len(answer)} ký tự")

    return {"question": question, "answer": answer, "retrieved_chunks": retrieved_chunks}

# API sinh caau hoir
@app.post("/Multiple_Choice_Questions")
async def Multiple_Choice_Questions(request: QuestionRequest):
    question = request.question
    retrieved_chunks = search_similar_chunks(question, top_k=3)

    if not retrieved_chunks:
        return {"question": question, "answer": "❌ Không tìm thấy tài liệu phù hợp."}

    context = "\n\n".join(f"{chunk['content']}" for chunk in retrieved_chunks)
    answer = generate_Multiple_Choice_Questions(question, context)

    return {"question": question, "answer": answer, "retrieved_chunks": retrieved_chunks}

@app.post("/Generate_Slide")
async def Generate_Slide_endpoint(request: QuestionRequest):
    question = request.question
    retrieved_chunks = search_similar_chunks(question, top_k=3)

    if not retrieved_chunks:
        return {"question": question, "answer": "❌ Không tìm thấy tài liệu phù hợp."}

    context = "\n\n".join(f"{chunk['content']}" for chunk in retrieved_chunks)
    
    # Lấy nội dung slide dạng text
    text_response = generate_Slide(question, context)
    
    # Parse text thành JSON
    try:
        import json
        import re
        
        # Xử lý text để lấy phần JSON
        json_text = re.sub(r'```json\s*|\s*```', '', text_response)
        slides_data = json.loads(json_text.strip())
        
        # Tạo PowerPoint
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
            import io
            
            # Tạo presentation
            prs = Presentation()
            
            # Xử lý từng slide
            for i, slide_info in enumerate(slides_data):
                try:
                    # Tạo slide với layout Title and Content
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    
                    # --- Tiêu đề slide ---
                    title_shape = slide.shapes.title
                    title_shape.text = slide_info.get("title", f"Slide {i+1}")
                    title_run = title_shape.text_frame.paragraphs[0].runs[0]
                    title_run.font.size = Pt(32)
                    title_run.font.bold = True
                    title_run.font.name = "Arial"
                    
                    # --- Nội dung ---
                    content_box = slide.placeholders[1]
                    content_frame = content_box.text_frame
                    content_frame.clear()
                    content_frame.word_wrap = True
                    
                    # Thêm từng section
                    for section in slide_info.get("sections", []):
                        # Heading
                        p_heading = content_frame.add_paragraph()
                        p_heading.text = section.get("heading", "")
                        p_heading.font.size = Pt(20)
                        p_heading.font.bold = True
                        p_heading.font.name = "Arial"
                        p_heading.alignment = PP_ALIGN.LEFT
                        p_heading.space_after = Pt(2)
                        
                        # Nội dung
                        content_text = section.get("content", "")
                        for line in content_text.strip().split('\n'):
                            p_content = content_frame.add_paragraph()
                            p_content.text = line.strip()
                            p_content.font.size = Pt(18)
                            p_content.font.name = "Arial"
                            p_content.alignment = PP_ALIGN.LEFT
                            p_content.space_after = Pt(6)
                            if line.strip().startswith("-"):
                                p_content.level = 1
                    
                    # --- Ghi chú ---
                    if "notes" in slide_info:
                        notes = slide.notes_slide.notes_text_frame
                        notes.clear()
                        p = notes.add_paragraph()
                        p.text = slide_info.get("notes", "")
                        p.font.size = Pt(14)
                        p.font.italic = True
                except Exception as slide_error:
                    print(f"Error processing slide {i+1}: {slide_error}")
            
            # Lưu file PowerPoint vào buffer
            ppt_buffer = io.BytesIO()
            prs.save(ppt_buffer)
            ppt_buffer.seek(0)
            
            # Encode PowerPoint buffer thành base64
            import base64
            ppt_base64 = base64.b64encode(ppt_buffer.read()).decode('utf-8')
            
            # Trả về cả JSON và PowerPoint base64
            return {
                "question": question, 
                "answer": text_response, 
                "retrieved_chunks": retrieved_chunks,
                "slides_data": slides_data,
                "ppt_base64": ppt_base64
            }
            
        except Exception as ppt_error:
            print(f"Error creating PowerPoint: {ppt_error}")
            import traceback
            traceback.print_exc()
            # Trả về JSON nếu có lỗi khi tạo PowerPoint
            return {
                "question": question, 
                "answer": text_response, 
                "retrieved_chunks": retrieved_chunks,
                "slides_data": slides_data,
                "error": f"Lỗi khi tạo PowerPoint: {str(ppt_error)}"
            }
            
    except Exception as json_error:
        print(f"Error parsing JSON: {json_error}")
        # Trả về text gốc nếu không parse được JSON
        return {
            "question": question, 
            "answer": text_response, 
            "retrieved_chunks": retrieved_chunks,
            "error": f"Lỗi khi parse JSON: {str(json_error)}"
        }

@app.post("/chat_topic")
async def chat_topic(request: QuestionRequest):
    try:
        context = request.question  # Đảm bảo nhận đúng dữ liệu từ request
        answer = get_chat_topic(context)
        return {"question": context, "answer": answer}
    except Exception as e:
        return {"error": str(e)} 
    
@app.post("/Quiz_Feedback")
async def Quiz_Feedback(request: QuestionRequest):
    try:
        context = request.question  # Đảm bảo nhận đúng dữ liệu từ request
        answer = get_Quiz_Feedback(context)
        return {"question": context, "answer": answer}
    except Exception as e:
        return {"error": str(e)}

@app.post("/Adaptive_Questions_endpoint")
async def Adaptive_Questions_endpoint(request: QuestionRequest):
    try:
        # Kiểm tra xem có previousPerformance trong request không
        if hasattr(request, 'previousPerformance'):
            previous_results_json = request.previousPerformance
        else:
            # Giả định thông tin hiệu suất trước đó có thể được gửi trong question
            previous_results_json = request.question
            
        topic = request.question  # Chủ đề được gửi qua trường question
        
        # Nếu request có định dạng đặc biệt, thử phân tích
        if hasattr(request, 'previous_performance'):
            previous_results_json = request.previous_performance
            
        answer = Adaptive_Questions(topic, previous_results_json)
        return {"question": topic, "answer": answer}
    except Exception as e:
        return {"error": str(e)}
    
@app.post("/grade_math_paper")
async def grade_math_paper(request: GradeRequest):
    try:
        student_image_path = request.student_image_path
        answer_key = request.answer_key 
        answer = Grade_math_paper(student_image_path,answer_key)
        return {"student_image_path": student_image_path, "answer": answer}
    except Exception as e:
        return {"error": str(e)}

@app.post("/search-test")
async def search_test(request: QuestionRequest):
    """API đơn giản để kiểm tra khả năng tìm kiếm các embedding"""
    try:
        question = request.question
        retrieved_chunks = search_similar_chunks(question, top_k=5)
        
        # In thông tin debugging
        print(f"Câu hỏi: {question}")
        print(f"Tìm thấy {len(retrieved_chunks)} chunk liên quan")
        
        # Lấy thông tin về collection
        collection_info = collection.get(include=["metadatas"])
        total_chunks = len(collection_info["ids"]) if "ids" in collection_info else 0
        
        return {
            "question": question, 
            "chunks_found": len(retrieved_chunks),
            "total_chunks_in_db": total_chunks,
            "chunks": retrieved_chunks
        }
    except Exception as e:
        print(f"⚠️ Lỗi khi thực hiện search-test: {str(e)}")
        return {"error": str(e)}

# Đảm bảo thư mục uploads tồn tại
UPLOAD_FOLDER = "uploads"
MATH_PAPERS_FOLDER = os.path.join(UPLOAD_FOLDER, "math_papers")

# Đảm bảo thư mục tồn tại khi khởi động ứng dụng
try:
    if not os.path.exists(UPLOAD_FOLDER):
        os.makedirs(UPLOAD_FOLDER)
        print(f"Đã tạo thư mục gốc: {UPLOAD_FOLDER}")
    
    if not os.path.exists(MATH_PAPERS_FOLDER):
        os.makedirs(MATH_PAPERS_FOLDER)
        print(f"Đã tạo thư mục lưu ảnh: {MATH_PAPERS_FOLDER}")
    else:
        print(f"Thư mục đã tồn tại: {MATH_PAPERS_FOLDER}")
        # Đếm số file trong thư mục
        file_count = len([f for f in os.listdir(MATH_PAPERS_FOLDER) if os.path.isfile(os.path.join(MATH_PAPERS_FOLDER, f))])
        print(f"Số lượng file hiện có: {file_count}")
except Exception as init_error:
    print(f"⚠️ Lỗi khi tạo thư mục ban đầu: {init_error}")

@app.post("/upload_image")
async def upload_image(file: UploadFile = File(...)):
    """
    API endpoint để upload ảnh và lưu vào thư mục cố định
    """
    try:
        # Tạo tên file duy nhất
        file_extension = file.filename.split(".")[-1] if "." in file.filename else "jpg"
        unique_filename = f"{uuid.uuid4()}.{file_extension}"
        
        # Sử dụng thư mục cố định không phân chia theo ngày
        relative_folder = MATH_PAPERS_FOLDER
        
        # Đảm bảo thư mục tồn tại với đầy đủ quyền truy cập
        try:
            if not os.path.exists(relative_folder):
                print(f"Tạo thư mục: {relative_folder}")
                os.makedirs(relative_folder, exist_ok=True)
                # Đợi để đảm bảo thư mục được tạo đầy đủ
                time.sleep(1)
                
            # Kiểm tra lại thư mục đã tồn tại chưa
            if not os.path.exists(relative_folder):
                print(f"Không thể tạo thư mục: {relative_folder}")
                return {"success": False, "error": f"Không thể tạo thư mục: {relative_folder}"}
        except Exception as folder_error:
            print(f"Lỗi khi tạo/kiểm tra thư mục: {folder_error}")
            return {"success": False, "error": f"Lỗi khi tạo thư mục: {str(folder_error)}"}
        
        # Đường dẫn tương đối của file
        relative_path = os.path.join(relative_folder, unique_filename)
        
        # Lưu file với xử lý lỗi chi tiết
        content = await file.read()
        if not content:
            return {"success": False, "error": "File rỗng hoặc không đọc được nội dung"}
            
        # Lưu file
        with open(relative_path, "wb") as buffer:
            buffer.write(content)
            
        # Kiểm tra file đã được lưu thành công chưa
        if not os.path.exists(relative_path) or os.path.getsize(relative_path) == 0:
            time.sleep(1)  # Đợi một chút cho hệ thống file
            if not os.path.exists(relative_path) or os.path.getsize(relative_path) == 0:
                return {"success": False, "error": "File không được lưu thành công"}
        
        print(f"Đã lưu file thành công tại: {relative_path}")
        # Trả về đường dẫn tương đối để API_Rag.py có thể truy cập
        return {"success": True, "file_path": relative_path}
    except Exception as e:
        print(f"Lỗi khi upload_image: {str(e)}")
        return {"success": False, "error": str(e)}

@app.post("/delete_image")
async def delete_image(file_data: dict):
    """
    API endpoint để xóa file ảnh từ thư mục uploads
    """
    try:
        file_path = file_data.get("file_path")
        if not file_path:
            return {"success": False, "error": "Thiếu đường dẫn file"}
        
        # Kiểm tra xem đường dẫn có hợp lệ không
        if not os.path.exists(file_path):
            return {"success": False, "error": f"File không tồn tại: {file_path}"}
        
        # Kiểm tra xem file có thuộc thư mục uploads không
        if "uploads" not in file_path:
            return {"success": False, "error": "Không được phép xóa file bên ngoài thư mục uploads"}
        
        # Xóa file
        os.remove(file_path)
        return {"success": True, "message": f"Đã xóa file: {file_path}"}
    except Exception as e:
        print(f"⚠️ Lỗi khi xóa file: {e}")
        return {"success": False, "error": str(e)}

@app.post("/verify_file")
async def verify_file(file_data: dict):
    """
    API endpoint để kiểm tra xem file tồn tại và có nội dung không
    """
    try:
        file_path = file_data.get("file_path")
        if not file_path:
            return {"success": False, "error": "Thiếu đường dẫn file"}
        
        # Chuẩn hóa đường dẫn
        file_path = os.path.normpath(file_path)
        
        # Kiểm tra file tồn tại
        if not os.path.exists(file_path):
            return {"success": False, "exists": False, "error": f"File không tồn tại: {file_path}"}
        
        # Kiểm tra file có dữ liệu
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            return {"success": False, "exists": True, "size": 0, "error": "File rỗng (0 byte)"}
        
        # Trả về thông tin chi tiết về file
        return {
            "success": True,
            "exists": True,
            "size": file_size,
            "path": file_path,
            "last_modified": os.path.getmtime(file_path)
        }
    except Exception as e:
        print(f"⚠️ Lỗi khi xác minh file: {e}")
        return {"success": False, "error": str(e)}

@app.post("/export_excel")
async def export_excel(request: ExcelExportRequest):
    """
    API endpoint để tạo file Excel từ kết quả chấm bài
    """
    try:
        print(f"Nhận yêu cầu xuất Excel với {len(request.results)} kết quả")
        
        # Tạo thư mục tạm để lưu file Excel
        temp_dir = os.path.join(UPLOAD_FOLDER, "temp_excel")
        os.makedirs(temp_dir, exist_ok=True)
        
        # Tạo workbook mới
        wb = Workbook()
        ws = wb.active
        ws.title = "Kết quả chấm điểm"
        
        # Thêm tiêu đề các cột
        headers = ["STT", "Họ và tên", "Lớp", "Điểm", "Nhận xét chi tiết", "Hình ảnh bài làm"]
        for col_num, header in enumerate(headers, 1):
            col_letter = get_column_letter(col_num)
            ws[f"{col_letter}1"] = header
            ws.column_dimensions[col_letter].width = 15  # Độ rộng cột mặc định
        
        # Thiết lập độ rộng cột cụ thể
        ws.column_dimensions['A'].width = 5   # STT
        ws.column_dimensions['B'].width = 25  # Họ và tên
        ws.column_dimensions['C'].width = 10  # Lớp
        ws.column_dimensions['D'].width = 10  # Điểm
        ws.column_dimensions['E'].width = 100 # Nhận xét
        ws.column_dimensions['F'].width = 50  # Hình ảnh bài làm
        
        # Thêm dữ liệu vào bảng
        for idx, item in enumerate(request.results, 1):
            row_num = idx + 1  # +1 vì hàng 1 là header
            
            try:
                # Điền dữ liệu cơ bản
                ws[f"A{row_num}"] = idx  # STT
                ws[f"B{row_num}"] = item.get('studentName', 'Không xác định')
                ws[f"C{row_num}"] = item.get('studentClass', 'Không xác định')
                ws[f"D{row_num}"] = item.get('totalScore', 'Không xác định')
                
                # Thêm nhận xét chi tiết - cắt bớt nếu quá dài
                full_result = item.get('fullResult', 'Không có dữ liệu')
                if len(full_result) > 32700:  # Giới hạn ký tự của Excel
                    full_result = full_result[:32700] + "... (đã cắt bớt)"
                ws[f"E{row_num}"] = full_result
                
                # Thêm hình ảnh bài làm
                image_path = item.get('imagePath', '')
                if image_path:
                    # Chuyển đường dẫn tương đối thành đường dẫn đầy đủ
                    full_image_path = os.path.join(MATH_PAPERS_FOLDER, os.path.basename(image_path))
                    if os.path.exists(full_image_path):
                        try:
                            # Tạo đối tượng Image từ file
                            img = XLImage(full_image_path)
                            
                            # Điều chỉnh kích thước ảnh nếu cần
                            # Giữ tỷ lệ khung hình nhưng giới hạn kích thước tối đa
                            max_width = 200
                            max_height = 200
                            
                            # Tính toán kích thước mới giữ nguyên tỷ lệ
                            if img.width > max_width or img.height > max_height:
                                ratio = min(max_width/img.width, max_height/img.height)
                                img.width = int(img.width * ratio)
                                img.height = int(img.height * ratio)
                            
                            # Chèn ảnh vào ô F
                            ws.add_image(img, f"F{row_num}")
                            
                            # Điều chỉnh chiều cao hàng để hiển thị ảnh
                            ws.row_dimensions[row_num].height = max(75, img.height * 0.75)  # 0.75 là hệ số chuyển đổi
                            
                        except Exception as img_error:
                            print(f"Lỗi khi chèn ảnh: {img_error}")
                            ws[f"F{row_num}"] = f"Lỗi khi chèn ảnh: {str(img_error)}"
                    else:
                        ws[f"F{row_num}"] = "Không tìm thấy hình ảnh"
                else:
                    ws[f"F{row_num}"] = "Không có hình ảnh"
                
                # Thiết lập wrap text cho cột nhận xét
                cell = ws[f"E{row_num}"]
                cell.alignment = openpyxl.styles.Alignment(wrap_text=True)
                
            except Exception as item_error:
                print(f"Lỗi khi thêm dòng {idx}: {item_error}")
                ws[f"E{row_num}"] = f"Lỗi khi xử lý dữ liệu: {str(item_error)}"
        
        # Tạo tên file Excel
        excel_filename = f"ket_qua_cham_diem_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        excel_path = os.path.join(temp_dir, excel_filename)
        
        # Lưu workbook
        try:
            wb.save(excel_path)
            print(f"Đã lưu file Excel tại: {excel_path}")
        except Exception as save_error:
            print(f"Lỗi khi lưu file Excel: {save_error}")
            return {"success": False, "error": f"Lỗi khi lưu file: {str(save_error)}"}
        
        # Trả về file Excel
        return FileResponse(
            path=excel_path,
            filename=excel_filename,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )
        
    except Exception as e:
        print(f"⚠️ Lỗi khi tạo file Excel: {e}")
        return {"success": False, "error": str(e)}

# Chạy serve
if __name__ == "__main__":
    import uvicorn
    print("🔌 API khởi chạy trên 0.0.0.0:8000")
    uvicorn.run(app, host="0.0.0.0", port=8000, reload=True)
